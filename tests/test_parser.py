from __future__ import annotations

import pytest

from app.parser import (
    DocxParser,
    PDFParser,
    PARSER_REGISTRY,
    PPTXParser,
    TEXT_EXTENSIONS,
    TextParser,
)


class TestTextParser:
    def test_parse_utf8(self):
        parser = TextParser()
        result = parser.parse(b"hello world", ".txt")
        assert result == "hello world"

    def test_parse_utf8_bom(self):
        parser = TextParser()
        result = parser.parse(b"\xef\xbb\xbfhello", ".md")
        assert result == "\ufeffhello"

    def test_parse_gb18030(self):
        parser = TextParser()
        result = parser.parse(b"\xc4\xe3\xba\xc3", ".txt")
        assert result == "\u4f60\u597d"

    def test_parse_latin1_fallback(self):
        parser = TextParser()
        result = parser.parse(b"\xff", ".txt")
        assert result == "\xff"

    def test_parse_empty_bytes(self):
        parser = TextParser()
        result = parser.parse(b"", ".txt")
        assert result == ""

    def test_parse_all_bytes_decodable_via_latin1(self):
        parser = TextParser()
        for byte_val in range(256):
            parser.parse(bytes([byte_val]), ".txt")


class TestPDFParser:
    def test_parse(self, pdf_bytes):
        parser = PDFParser()
        result = parser.parse(pdf_bytes, ".pdf")
        assert "Hello, PDF!" in result

    def test_parse_empty(self):
        import fitz
        import io

        doc = fitz.open()
        doc.new_page()
        buf = io.BytesIO()
        doc.save(buf)
        doc.close()
        parser = PDFParser()
        result = parser.parse(buf.getvalue(), ".pdf")
        assert result == ""

    def test_parse_multi_page(self):
        import fitz
        import io

        doc = fitz.open()
        page1 = doc.new_page()
        page1.insert_text((50, 50), "Page One")
        page2 = doc.new_page()
        page2.insert_text((50, 50), "Page Two")
        buf = io.BytesIO()
        doc.save(buf)
        doc.close()
        parser = PDFParser()
        result = parser.parse(buf.getvalue(), ".pdf")
        assert "Page One" in result
        assert "Page Two" in result


class TestDocxParser:
    def test_parse(self, docx_bytes):
        parser = DocxParser()
        result = parser.parse(docx_bytes, ".docx")
        assert "Hello, DOCX!" in result

    def test_parse_empty(self):
        import io
        from docx import Document

        doc = Document()
        buf = io.BytesIO()
        doc.save(buf)
        parser = DocxParser()
        result = parser.parse(buf.getvalue(), ".docx")
        assert result == ""

    def test_parse_multiple_paragraphs(self):
        import io
        from docx import Document

        doc = Document()
        doc.add_paragraph("First para")
        doc.add_paragraph("Second para")
        buf = io.BytesIO()
        doc.save(buf)
        parser = DocxParser()
        result = parser.parse(buf.getvalue(), ".docx")
        assert "First para" in result
        assert "Second para" in result


class TestPPTXParser:
    def test_parse(self, pptx_bytes):
        parser = PPTXParser()
        result = parser.parse(pptx_bytes, ".pptx")
        assert "Hello, PPTX!" in result

    def test_parse_empty(self):
        import io
        from pptx import Presentation

        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)
        parser = PPTXParser()
        result = parser.parse(buf.getvalue(), ".pptx")
        assert result == ""

    def test_parse_with_table(self):
        import io
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(1))
        table = table_shape.table
        table.cell(0, 0).text = "A1"
        table.cell(0, 1).text = "B1"
        table.cell(1, 0).text = "A2"
        table.cell(1, 1).text = "B2"
        buf = io.BytesIO()
        prs.save(buf)
        parser = PPTXParser()
        result = parser.parse(buf.getvalue(), ".pptx")
        assert "A1" in result
        assert "B1" in result
        assert "A2" in result
        assert "B2" in result


class TestCorruptFiles:
    def test_corrupt_pdf_raises_value_error(self):
        parser = PDFParser()
        corrupt_bytes = b"not a real pdf content\x00\xff\xfe"
        with pytest.raises(ValueError, match="corrupted|Unable to open"):
            parser.parse(corrupt_bytes, ".pdf")

    def test_corrupt_docx_raises_value_error(self):
        parser = DocxParser()
        corrupt_bytes = b"not a real docx\x00\xff\xfe"
        with pytest.raises(ValueError, match="corrupted|Unable to open"):
            parser.parse(corrupt_bytes, ".docx")

    def test_corrupt_pptx_raises_value_error(self):
        parser = PPTXParser()
        corrupt_bytes = b"not a real pptx\x00\xff\xfe"
        with pytest.raises(ValueError, match="corrupted|Unable to open"):
            parser.parse(corrupt_bytes, ".pptx")

    def test_empty_pdf_raises_no_error(self):
        import fitz
        import io

        doc = fitz.open()
        doc.new_page()
        buf = io.BytesIO()
        doc.save(buf)
        doc.close()
        parser = PDFParser()
        result = parser.parse(buf.getvalue(), ".pdf")
        assert result == ""


class TestParserRegistry:
    def test_contains_text_extensions(self):
        for ext in TEXT_EXTENSIONS:
            assert ext in PARSER_REGISTRY, f"Missing extension: {ext}"

    def test_contains_binary_extensions(self):
        for ext in (".pdf", ".docx", ".pptx"):
            assert ext in PARSER_REGISTRY, f"Missing extension: {ext}"

    def test_text_parser_is_default_for_text_extensions(self):
        for ext in TEXT_EXTENSIONS:
            assert isinstance(PARSER_REGISTRY[ext], TextParser), f"Wrong parser for {ext}"

    def test_pdf_parser(self):
        assert isinstance(PARSER_REGISTRY[".pdf"], PDFParser)

    def test_docx_parser(self):
        assert isinstance(PARSER_REGISTRY[".docx"], DocxParser)

    def test_pptx_parser(self):
        assert isinstance(PARSER_REGISTRY[".pptx"], PPTXParser)
