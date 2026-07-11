from __future__ import annotations

import io
from pathlib import Path
from typing import Any

import pytest

from core.ingest import IngestionEngine
from core.project import ProjectConfig, ProjectContext
from core.retriever import HybridRetriever
from core.runtime import WorkspaceRuntime


# ── Fake embedding model ──────────────────────────────────────────

class FakeEmbeddingModel:
    """Deterministic embedding for tests: checks keyword presence in each dim."""

    def get_text_embedding_batch(self, texts: list[str]) -> list[list[float]]:
        return [self.get_text_embedding(text) for text in texts]

    def get_text_embedding(self, text: str) -> list[float]:
        text = text.lower()
        return [
            1.0 if "retrieval" in text or "检索" in text else 0.0,
            1.0 if "generator" in text or "生成" in text else 0.0,
            1.0 if "config" in text or "配置" in text else 0.0,
        ]

    def get_query_embedding(self, query: str) -> list[float]:
        return self.get_text_embedding(query)


# ── Mock LLM infrastructure ───────────────────────────────────────

class _Message:
    def __init__(self, content: str) -> None:
        self.content = content


class _Delta:
    def __init__(self, content: str) -> None:
        self.content = content


class _Choice:
    """Mock a single chat completion choice."""

    def __init__(self, content: str, for_stream: bool = False) -> None:
        if for_stream:
            self.delta = _Delta(content)
        else:
            self.message = _Message(content)


class _Response:
    """Mock a non-streaming response."""

    def __init__(self, content: str) -> None:
        self.choices = [_Choice(content)]


class _Chunk:
    """Mock a streaming chunk — only delta, no message."""

    def __init__(self, content: str) -> None:
        self.choices = [_Choice(content, for_stream=True)]


class _Completions:
    """Mock llm.chat.completions — matches openai's chat.completions.create()."""

    def __init__(self, sync_response: str, stream_chunks: list[_Chunk] | None = None) -> None:
        self._sync_response = sync_response
        self._stream_chunks = stream_chunks

    def create(
        self,
        model: str | None = None,
        messages: list[dict[str, str]] | None = None,
        temperature: float | None = None,
        max_tokens: int | None = None,
        stream: bool = False,
    ) -> Any:
        if stream:
            return iter(self._stream_chunks or [_Chunk("")])
        return _Response(self._sync_response)


class _Chat:
    """Mock llm.chat — provides .completions.create() matching openai SDK path."""

    def __init__(self, sync_response: str = "", stream_chunks: list[_Chunk] | None = None) -> None:
        self.completions = _Completions(sync_response, stream_chunks)


def _make_chat(sync: str, stream: list[_Chunk] | None = None) -> _Chat:
    return _Chat(sync_response=sync, stream_chunks=stream or [_Chunk(sync)])


class FakeLLM:
    """LLM mock that returns a fixed answer (non-streaming and streaming)."""

    def __init__(self) -> None:
        self.chat = _make_chat(
            sync="回答基于检索上下文。[1]\n\n引用来源：\n[1] unknown",
            stream=[_Chunk("回答基于检索上下文。[1]\n\n引用来源：\n[1] unknown")],
        )


class EmptyLLM:
    """LLM mock that returns empty content."""

    def __init__(self) -> None:
        self.chat = _make_chat(sync="")


class StreamingLLM:
    """LLM mock that yields multiple stream chunks."""

    def __init__(self) -> None:
        self.chat = _make_chat(
            sync="第一段第二段",
            stream=[_Chunk("第一段"), _Chunk("第二段")],
        )


# ── Test helpers ──────────────────────────────────────────────────

def make_project_config(tmp_path: Path, **overrides: Any) -> ProjectConfig:
    """Create a ProjectConfig in a temp directory with sensible test defaults.

    Sets up a minimal .lore/vectordb/ structure so core module engines
    can operate.  Files should be written directly under *tmp_path*.
    """
    co_dir = tmp_path / ".lore"
    (co_dir / "vectordb").mkdir(parents=True, exist_ok=True)

    config = ProjectConfig.load(co_dir / "config.toml")
    config.chunk_size = 400
    config.chunk_overlap = 20
    config.top_k = 5
    config.retrieval_candidate_k = 10
    config.similarity_cutoff = 0.0
    for key, value in overrides.items():
        if hasattr(config, key):
            setattr(config, key, value)
    return config


def build_retrieval_results(
    tmp_path: Path,
    query: str = "retrieval",
) -> tuple[ProjectConfig, Any]:
    """Ingest two test documents under *tmp_path* and return (config, retrieval_results).

    Uses ``core.ingest.IngestionEngine`` and ``core.retriever.HybridRetriever``
    so the test seam matches production.
    """
    config = make_project_config(tmp_path)
    embedding_model = FakeEmbeddingModel()
    engine = IngestionEngine(config=config, root=tmp_path, embedding_model=embedding_model)

    (tmp_path / "retrieval.md").write_text(
        "Hybrid retrieval combines vector retrieval with BM25 for better recall.",
        encoding="utf-8",
    )
    (tmp_path / "generator.md").write_text(
        "The generator builds answers from retrieved context and cites sources.",
        encoding="utf-8",
    )

    engine.add_files(engine.scan_files())
    retriever = HybridRetriever(config=config, vector_store=engine.vector_store, embedding_model=embedding_model)
    results = retriever.retrieve(query, mode="hybrid")
    return config, results


# ── WorkspaceRuntime helper ────────────────────────────────────────


def make_runtime(tmp_path: Path) -> WorkspaceRuntime:
    """创建带假引擎的 WorkspaceRuntime，用于测试。"""
    co_dir = tmp_path / ".lore"
    (co_dir / "vectordb").mkdir(parents=True, exist_ok=True)

    ctx = ProjectContext(tmp_path)
    ctx.config.chunk_size = 200
    ctx.config.chunk_overlap = 20
    ctx.config.top_k = 5
    ctx.config.retrieval_candidate_k = 10
    ctx.config.similarity_cutoff = 0.0
    ctx.embedding_model = FakeEmbeddingModel()
    ctx.llm = FakeLLM()
    ctx.setup_engines()
    return WorkspaceRuntime._from_ctx(ctx)


# ── Document parsing fixtures ─────────────────────────────────────

def make_pdf_bytes(text: str = "Hello, PDF!") -> bytes:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), text)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def make_docx_bytes(text: str = "Hello, DOCX!") -> bytes:
    from docx import Document

    doc = Document()
    doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_pptx_bytes(text: str = "Hello, PPTX!") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2)).text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def pdf_bytes() -> bytes:
    return make_pdf_bytes()


@pytest.fixture
def docx_bytes() -> bytes:
    return make_docx_bytes()


@pytest.fixture
def pptx_bytes() -> bytes:
    return make_pptx_bytes()
