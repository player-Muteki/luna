from __future__ import annotations

import io
import logging
from abc import ABC, abstractmethod
from pathlib import Path

logger = logging.getLogger(__name__)

# Maximum file size (in bytes) for binary format parsers (PDF, DOCX, PPTX).
# Files larger than this threshold raise a descriptive error before attempting
# to parse, preventing OOM on extremely large documents.
MAX_BINARY_FILE_BYTES: int = 200 * 1024 * 1024  # 200 MiB


class DocumentParser(ABC):
    @abstractmethod
    def parse(self, file_bytes: bytes, file_ext: str) -> str:
        ...


class TextParser(DocumentParser):
    def parse(self, file_bytes: bytes, file_ext: str) -> str:
        encodings = ("utf-8", "utf-8-sig", "gb18030", "latin-1")
        for encoding in encodings:
            try:
                return file_bytes.decode(encoding)
            except UnicodeDecodeError:
                continue
        raise ValueError(f"Unable to decode file with any of {encodings}")


class PDFParser(DocumentParser):
    def parse(self, file_bytes: bytes, file_ext: str) -> str:
        if len(file_bytes) > MAX_BINARY_FILE_BYTES:
            raise ValueError(
                f"PDF file exceeds size limit of {MAX_BINARY_FILE_BYTES // (1024*1024)} MiB"
            )
        try:
            import fitz
        except ImportError as exc:
            raise RuntimeError("pymupdf (fitz) is required to parse PDF files") from exc

        try:
            doc = fitz.open(stream=file_bytes, filetype="pdf")
        except Exception as exc:
            raise ValueError(f"Unable to open PDF file (possibly corrupted): {exc}") from exc

        try:
            pages: list[str] = []
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                if text.strip():
                    pages.append(text)
            return "\n\n".join(pages)
        finally:
            doc.close()


class DocxParser(DocumentParser):
    def parse(self, file_bytes: bytes, file_ext: str) -> str:
        if len(file_bytes) > MAX_BINARY_FILE_BYTES:
            raise ValueError(
                f"DOCX file exceeds size limit of {MAX_BINARY_FILE_BYTES // (1024*1024)} MiB"
            )
        try:
            from docx import Document
        except ImportError as exc:
            raise RuntimeError("python-docx is required to parse .docx files") from exc

        try:
            doc = Document(io.BytesIO(file_bytes))
        except Exception as exc:
            raise ValueError(f"Unable to open DOCX file (possibly corrupted): {exc}") from exc

        paragraphs: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        return "\n\n".join(paragraphs)


class PPTXParser(DocumentParser):
    def parse(self, file_bytes: bytes, file_ext: str) -> str:
        if len(file_bytes) > MAX_BINARY_FILE_BYTES:
            raise ValueError(
                f"PPTX file exceeds size limit of {MAX_BINARY_FILE_BYTES // (1024*1024)} MiB"
            )
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("python-pptx is required to parse .pptx files") from exc

        try:
            prs = Presentation(io.BytesIO(file_bytes))
        except Exception as exc:
            raise ValueError(f"Unable to open PPTX file (possibly corrupted): {exc}") from exc

        slides: list[str] = []
        for slide in prs.slides:
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells]
                        slide_texts.append(" | ".join(row_cells))
            if slide_texts:
                slides.append("\n".join(slide_texts))
        return "\n\n".join(slides)


TEXT_EXTENSIONS = frozenset({
    # Programming languages
    ".c", ".cpp", ".cs", ".go", ".h", ".hpp", ".java", ".js", ".jsx",
    ".kt", ".kts", ".lua", ".php", ".pl", ".pm", ".py", ".rb", ".rs",
    ".scala", ".swift", ".ts", ".tsx",
    # Documentation
    ".md", ".mdx", ".rst", ".tex", ".sty", ".cls", ".bib", ".bst",
    ".html", ".htm", ".txt",
    # Data & config
    ".json", ".yaml", ".yml", ".toml", ".xml", ".ini", ".cfg", ".conf", ".env",
    ".csv", ".sql", ".log",
    # Shell & scripts
    ".sh", ".bash", ".zsh", ".m", ".mm", ".r", ".R",
    # Build & infra
    ".tf", ".cmake", ".gradle",
})


def _build_default_registry() -> dict[str, DocumentParser]:
    text_parser = TextParser()
    pdf_parser = PDFParser()
    docx_parser = DocxParser()
    pptx_parser = PPTXParser()

    registry: dict[str, DocumentParser] = {}
    for ext in TEXT_EXTENSIONS:
        registry[ext] = text_parser
    registry[".pdf"] = pdf_parser
    registry[".docx"] = docx_parser
    registry[".pptx"] = pptx_parser
    return registry


PARSER_REGISTRY: dict[str, DocumentParser] = _build_default_registry()
